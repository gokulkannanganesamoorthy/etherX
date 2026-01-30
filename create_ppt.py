from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title_text, content_text_list):
    slide_layout = prs.slide_layouts[1]  # Bullet list layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = title_text
    
    # Content
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = content_text_list[0] # First item
    
    for item in content_text_list[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "EtherX Sentinel: Advanced AI WAF"
    subtitle.text = "Team KADAVUL\nPresenter: Gokul Kannan Ganesamoorthy"

    # Slide 1: The Problem
    add_slide(prs, "The Problem: Why Legacy Firewalls Fail", [
        "Rule-Based Limitations: Traditional WAFs rely on static Regex signatures.",
        "The Zero-Day Threat: Attackers constantly mutate payloads to bypass filters.",
        "Maintenance Nightmare: Maintaining thousands of rules is inefficient.",
        "False Positives: Strict rules often block legitimate user traffic."
    ])

    # Slide 2: The Solution - EtherX Sentinel
    add_slide(prs, "The Solution: A Paradigm Shift", [
        "\"Don't look for the attack. Look for the intention.\"",
        "Semantic Understanding: Uses NLP to understand the meaning of requests.",
        "Behavioral Analysis: Trained on 'normal' traffic deviations.",
        "No Signatures: Detects 0-day XSS and SQLi attacks without Regex."
    ])

    # Slide 3: The Deep Learning Engine
    add_slide(prs, "The Deep Learning Engine", [
        "1. Ingestion: High-speed traffic interception via FastAPI.",
        "2. Embedding: Sentence Transformers convert payloads to 384-d vectors.",
        "3. Neural Encoder: PyTorch Autoencoder reconstructs these vectors.",
        "4. Anomaly Detection: High reconstruction error = Block request.",
        "Result: <10ms inference latency."
    ])

    # Slide 4: Technical Deep Dive
    add_slide(prs, "Technical Deep Dive: Under the Hood", [
        "1. The AI Model: PyTorch Autoencoder (Compression Network).",
        "2. Embedding: 'all-MiniLM-L6-v2' (384-d Vector Space).",
        "3. Training Data: 'benign_traffic.txt' (Learns Normality).",
        "4. Risk Analysis (JSON Brain):",
        "   - reconstruction_error: >0.02 = Anomaly (Zero-Day).",
        "   - neural_anomaly: Boolean flag for high-confidence blocks."
    ])

    # Slide 5: The Training Pipeline
    add_slide(prs, "The Training Pipeline", [
        "1. Data Gen: 'generate_traffic.py' creates synthetic 'valid' traffic.",
        "2. Vectorization: Traffic -> Embeddings (Lists of numbers).",
        "3. Self-Supervised Learning: Model forces itself to learn patterns.",
        "   - Loss Function: Mean Squared Error (MSE).",
        "   - Goal: Output ≈ Input.",
        "4. Result: Model becomes an expert at 'Normality'."
    ])

    # Slide 6: Key Innovations
    add_slide(prs, "Key Innovations", [
        "Holographic Dashboard: Real-time, WebSocket-powered Neural Grid UI.",
        "Live Threat Intel: Active monitoring of SQLi, XSS, and Anomalies.",
        "Persistent Memory: SQLite-backed behavioral logging.",
        "Cyberpunk Aesthetics: Designed for the modern SOC."
    ])

    # Slide 5: Architecture
    add_slide(prs, "Architecture", [
        "Client Traffic -> Ingestion Layer (FastAPI)",
        "-> AI Model (SentenceTransformer + Autoencoder)",
        "-> Decision (Allow/Block)",
        "-> Dashboard (WebSocket Stream)"
    ])

    # Slide 6: Live Demo
    add_slide(prs, "Live Demo Plan", [
        "1. Normal Traffic: Browsing the site (Allowed).",
        "2. SQL Injection: Attempting 'OR 1=1' (Blocked by AI).",
        "3. XSS Attack: Trying '<script>' tags (Blocked).",
        "4. Obfuscation: Complex encodings (Blocked by Semantic Analysis).",
        "5. Dashboard: Real-time visualization of these events."
    ])

    # Slide 7: Conclusion
    add_slide(prs, "Conclusion", [
        "EtherX Sentinel represents the future of adaptive application security.",
        "It learns, evolves, and protects without manual intervention.",
        "Team KADAVUL is proud to present this next-gen solution."
    ])
    
    # Q&A Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Q&A"
    
    prs.save('EtherX_Sentinel.pptx')
    print("Presentation saved as EtherX_Sentinel.pptx")

if __name__ == "__main__":
    create_presentation()
